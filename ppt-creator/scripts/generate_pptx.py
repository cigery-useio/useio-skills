#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
PPT Creator - python-pptx 生成脚本

使用 python-pptx 库生成高质量 PowerPoint 演示文稿。
支持创建新 PPT 和修改已有 PPT。

用法：
    python generate_pptx.py --input slides.json --output output.pptx --skill-dir /path/to/ppt-creator
    python generate_pptx.py --input slides.json --source existing.pptx --output modified.pptx --skill-dir /path/to/ppt-creator
    python generate_pptx.py --input slides.json --output output.pptx --skill-dir /path/to/ppt-creator --aspect-ratio 16:9 --template template.pptx

参数：
    --input        : slides JSON 文件路径（必填）
    --output       : 输出 PPTX 文件路径（必填）
    --skill-dir    : 技能目录绝对路径（优先于 __file__ 自探测，用于定位 vendor/libs）
    --source       : 修改模式时，源 PPTX 文件路径（可选）
    --aspect-ratio : 宽高比，16:9（默认）或 4:3（仅 create 模式生效，modify 保留源文件宽高比）
    --template     : 模板 PPTX 文件路径（仅 create 模式，基于模板创建保留母版/主题）
"""

import os
import sys
import json
import argparse
import base64
import io
import shutil
import glob
import time
from pathlib import Path

# 添加 vendor/libs 到 sys.path
# 优先使用 --skill-dir 参数，其次 __file__ 自探测
_arg_skill_dir = None
for i, arg in enumerate(sys.argv):
    if arg == "--skill-dir" and i + 1 < len(sys.argv):
        _arg_skill_dir = sys.argv[i + 1]
        break

if _arg_skill_dir:
    skill_dir = os.path.abspath(_arg_skill_dir)
else:
    script_dir = os.path.dirname(os.path.abspath(__file__))
    skill_dir = os.path.dirname(script_dir)

libs_path = os.path.join(skill_dir, "vendor", "libs")
if os.path.isdir(libs_path):
    sys.path.insert(0, libs_path)
else:
    print(f"错误：vendor/libs 目录不存在: {libs_path}", file=sys.stderr)
    print("请使用 generate_stdlib.py（路线 C）", file=sys.stderr)
    sys.exit(2)  # 退出码 2 区分「依赖缺失」和「代码错误」

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.chart import XL_CHART_TYPE
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False
    print("错误：python-pptx 未安装，请使用路线 C（纯标准库）", file=sys.stderr)
    sys.exit(2)


# ============================================================
# 主题配置
# ============================================================

DEFAULT_THEME = {
    "primaryColor": "2B579A",
    "fontFamily": "微软雅黑"
}

# 宽高比尺寸（EMU 转 Inches）
ASPECT_RATIOS = {
    "16:9": {"width": Inches(13.333), "height": Inches(7.5)},
    "4:3": {"width": Inches(10), "height": Inches(7.5)},
}

# 支持的布局类型
SUPPORTED_LAYOUTS = {"title", "content", "two_content", "table", "image", "section", "blank", "chart"}


def hex_to_rgb(hex_color):
    """将十六进制颜色转换为 RGBColor"""
    hex_color = hex_color.lstrip('#')
    if len(hex_color) == 6:
        return RGBColor(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16)
        )
    return RGBColor(0x2B, 0x57, 0x9A)  # 默认主题色


# ============================================================
# JSON 校验
# ============================================================

def validate_slides_json(data):
    """校验 slides JSON 格式

    校验顶层必填字段和每个 slide 的 layout 字段。
    校验失败时输出明确错误信息并 sys.exit(1)。
    """
    if not isinstance(data, dict):
        print("错误：JSON 顶层必须为对象", file=sys.stderr)
        sys.exit(1)

    if not isinstance(data.get("title"), str):
        print("错误：缺少必填字段 title 或类型不是 string", file=sys.stderr)
        sys.exit(1)

    slides = data.get("slides")
    if not isinstance(slides, list) or len(slides) == 0:
        print("错误：缺少必填字段 slides 或不是非空数组", file=sys.stderr)
        sys.exit(1)

    for i, slide in enumerate(slides):
        if not isinstance(slide, dict):
            print(f"错误：第 {i+1} 张幻灯片不是对象", file=sys.stderr)
            sys.exit(1)
        layout = slide.get("layout")
        if layout not in SUPPORTED_LAYOUTS:
            print(f"错误：第 {i+1} 张幻灯片 layout '{layout}' 不在支持列表内: {SUPPORTED_LAYOUTS}", file=sys.stderr)
            sys.exit(1)

    # 校验 aspectRatio（若提供）
    aspect_ratio = data.get("theme", {}).get("aspectRatio")
    if aspect_ratio and aspect_ratio not in ASPECT_RATIOS:
        print(f"错误：aspectRatio '{aspect_ratio}' 不在支持列表内: {list(ASPECT_RATIOS.keys())}", file=sys.stderr)
        sys.exit(1)


# ============================================================
# 修改前备份（模块 8）
# ============================================================

def get_app_root(skill_dir_path):
    """从 skill_dir 推算 app_root（向上两级）"""
    return os.path.dirname(os.path.dirname(skill_dir_path))


def backup_source(source_path, skill_dir_path):
    """修改前自动备份源文件到 {app_root}/data/ppt-backups/

    备份文件名格式：{原文件名}_{时间戳}.pptx
    备份失败不中断主流程，仅 stderr 警告。
    """
    if not os.path.isfile(source_path):
        return

    app_root = get_app_root(skill_dir_path)
    backup_dir = os.path.join(app_root, "data", "ppt-backups")

    try:
        os.makedirs(backup_dir, exist_ok=True)
    except OSError as e:
        print(f"警告：创建备份目录失败: {e}", file=sys.stderr)
        return

    # 生成备份文件名
    source_filename = os.path.basename(source_path)
    timestamp = time.strftime("%Y%m%d_%H%M%S")
    backup_filename = f"{source_filename}_{timestamp}.pptx"
    backup_path = os.path.join(backup_dir, backup_filename)

    try:
        shutil.copy2(source_path, backup_path)
        print(f"备份已创建: {backup_path}", file=sys.stderr)
    except OSError as e:
        print(f"警告：备份源文件失败: {e}", file=sys.stderr)
        return

    # 双层清理
    cleanup_old_backups(backup_dir, source_filename)


def cleanup_old_backups(backup_dir, source_filename, max_keep=5, max_total_mb=500):
    """双层清理：按源文件分组 + 总量控制

    第一层：同一源文件名保留最近 max_keep 个备份
    第二层：整个备份目录总大小超过 max_total_mb 时从最旧开始删除
    """
    # 第一层：按源文件分组清理
    pattern = os.path.join(backup_dir, f"{source_filename}_*.pptx")
    backups = sorted(glob.glob(pattern), key=os.path.getmtime, reverse=True)
    for old_backup in backups[max_keep:]:
        try:
            os.remove(old_backup)
        except OSError:
            pass

    # 第二层：总量控制
    all_backups = sorted(
        glob.glob(os.path.join(backup_dir, "*_*.pptx")),
        key=os.path.getmtime,
        reverse=True  # 最新的在前
    )
    try:
        total_size_mb = sum(os.path.getsize(f) for f in all_backups) / (1024 * 1024)
    except OSError:
        return

    if total_size_mb > max_total_mb:
        # 从最旧的开始删除（reversed），直至低于上限
        for old_file in reversed(all_backups):
            if total_size_mb <= max_total_mb:
                break
            try:
                file_size_mb = os.path.getsize(old_file) / (1024 * 1024)
                os.remove(old_file)
                total_size_mb -= file_size_mb
            except OSError:
                pass


# ============================================================
# 布局处理函数
# ============================================================

def add_title_slide(prs, slide_data, theme):
    """添加标题页"""
    slide = prs.slides.add_slide(prs.slide_layouts[0])  # 标题布局

    title_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "title"), None)
    subtitle_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "subtitle"), None)

    if title_elem and slide.shapes.title:
        slide.shapes.title.text = title_elem.get("text", "")
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(36)
            paragraph.font.bold = True
            paragraph.font.color.rgb = hex_to_rgb(theme.get("primaryColor", DEFAULT_THEME["primaryColor"]))
            paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    if subtitle_elem:
        if len(slide.placeholders) > 1:
            subtitle_shape = slide.placeholders[1]
            subtitle_shape.text = subtitle_elem.get("text", "")
            for paragraph in subtitle_shape.text_frame.paragraphs:
                paragraph.font.size = Pt(20)
                paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    return slide


def add_content_slide(prs, slide_data, theme):
    """添加内容页（标题 + 要点列表）"""
    slide = prs.slides.add_slide(prs.slide_layouts[1])  # 标题和内容布局

    title_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "title"), None)
    bullets_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "bullets"), None)

    if title_elem and slide.shapes.title:
        slide.shapes.title.text = title_elem.get("text", "")
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    if bullets_elem:
        if len(slide.placeholders) > 1:
            body_shape = slide.placeholders[1]
            text_frame = body_shape.text_frame
            text_frame.clear()

            items = bullets_elem.get("items", [])
            for i, item in enumerate(items):
                if i == 0:
                    p = text_frame.paragraphs[0]
                else:
                    p = text_frame.add_paragraph()
                p.text = item
                p.font.size = Pt(16)  # 16:9 适配：从 Pt(18) 调整为 Pt(16)
                p.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])
                p.space_after = Pt(10)  # 增加行间距
                p.level = 0

    return slide


def add_two_content_slide(prs, slide_data, theme):
    """添加双栏页"""
    slide = prs.slides.add_slide(prs.slide_layouts[3])  # 双栏布局

    title_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "title"), None)
    left_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "left_bullets"), None)
    right_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "right_bullets"), None)

    if title_elem and slide.shapes.title:
        slide.shapes.title.text = title_elem.get("text", "")
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    # 左栏
    if left_elem and len(slide.placeholders) > 1:
        body_shape = slide.placeholders[1]
        text_frame = body_shape.text_frame
        text_frame.clear()
        items = left_elem.get("items", [])
        for i, item in enumerate(items):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)  # 16:9 适配
            p.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    # 右栏
    if right_elem and len(slide.placeholders) > 2:
        body_shape = slide.placeholders[2]
        text_frame = body_shape.text_frame
        text_frame.clear()
        items = right_elem.get("items", [])
        for i, item in enumerate(items):
            p = text_frame.paragraphs[0] if i == 0 else text_frame.add_paragraph()
            p.text = item
            p.font.size = Pt(14)  # 16:9 适配
            p.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    return slide


def add_table_slide(prs, slide_data, theme):
    """添加表格页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 仅标题布局

    title_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "title"), None)
    table_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "table"), None)

    if title_elem and slide.shapes.title:
        slide.shapes.title.text = title_elem.get("text", "")
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    if table_elem:
        headers = table_elem.get("headers", [])
        rows = table_elem.get("rows", [])
        num_rows = len(rows) + 1
        num_cols = len(headers) if headers else (len(rows[0]) if rows else 1)

        # 添加表格 — 16:9 适配：宽度从 Inches(8) 调整为 Inches(11)
        left = Inches(1)
        top = Inches(2)
        width = Inches(11)
        height = Inches(min(0.3 * num_rows, 5.0))  # 限制最大高度防止溢出

        table_shape = slide.shapes.add_table(num_rows, num_cols, left, top, width, height)
        table = table_shape.table

        # 填充表头
        for j, header in enumerate(headers):
            cell = table.cell(0, j)
            cell.text = header
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.bold = True
                paragraph.font.size = Pt(14)
                paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])
            cell.fill.solid()
            cell.fill.fore_color.rgb = hex_to_rgb(theme.get("primaryColor", DEFAULT_THEME["primaryColor"]))
            for paragraph in cell.text_frame.paragraphs:
                paragraph.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

        # 填充数据行
        for i, row in enumerate(rows):
            for j, cell_text in enumerate(row):
                if j < num_cols:
                    cell = table.cell(i + 1, j)
                    cell.text = str(cell_text)
                    for paragraph in cell.text_frame.paragraphs:
                        paragraph.font.size = Pt(12)
                        paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    return slide


def add_image_slide(prs, slide_data, theme):
    """添加图片页

    支持两种图片来源：
    - data: base64 编码的图片数据
    - filePath: 文件路径（与 data 二选一）
    """
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 仅标题布局

    title_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "title"), None)
    image_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "image"), None)

    if title_elem and slide.shapes.title:
        slide.shapes.title.text = title_elem.get("text", "")
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    if image_elem:
        image_stream = None
        should_close = False
        # 优先使用 filePath，其次 data
        file_path = image_elem.get("filePath")
        image_data = image_elem.get("data", "")

        if file_path:
            try:
                image_stream = open(file_path, 'rb')
                should_close = True
            except Exception as e:
                print(f"警告：读取图片文件失败: {e}", file=sys.stderr)
        elif image_data:
            try:
                image_bytes = base64.b64decode(image_data)
                image_stream = io.BytesIO(image_bytes)
            except Exception as e:
                print(f"警告：图片 base64 解码失败: {e}", file=sys.stderr)

        if image_stream:
            try:
                width = Inches(image_elem.get("width", 800) / 96)  # 像素转英寸
                height = Inches(image_elem.get("height", 600) / 96)
                slide.shapes.add_picture(image_stream, Inches(1), Inches(2), width, height)
            except Exception as e:
                print(f"警告：图片添加失败: {e}", file=sys.stderr)
            finally:
                if should_close:
                    image_stream.close()

    return slide


def add_section_slide(prs, slide_data, theme):
    """添加章节分隔页"""
    slide = prs.slides.add_slide(prs.slide_layouts[2])  # 节标题布局

    title_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "title"), None)
    subtitle_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "subtitle"), None)

    if title_elem and slide.shapes.title:
        slide.shapes.title.text = title_elem.get("text", "")
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.size = Pt(40)
            paragraph.font.bold = True
            paragraph.font.color.rgb = hex_to_rgb(theme.get("primaryColor", DEFAULT_THEME["primaryColor"]))
            paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    if subtitle_elem and len(slide.placeholders) > 1:
        subtitle_shape = slide.placeholders[1]
        subtitle_shape.text = subtitle_elem.get("text", "")
        for paragraph in subtitle_shape.text_frame.paragraphs:
            paragraph.font.size = Pt(24)
            paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    return slide


def add_blank_slide(prs, slide_data, theme):
    """添加空白页"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 空白布局
    return slide


def add_chart_slide(prs, slide_data, theme):
    """添加图表页"""
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # 仅标题布局

    title_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "title"), None)
    chart_elem = next((e for e in slide_data.get("elements", []) if e.get("type") == "chart"), None)

    if title_elem and slide.shapes.title:
        slide.shapes.title.text = title_elem.get("text", "")
        for paragraph in slide.shapes.title.text_frame.paragraphs:
            paragraph.font.name = theme.get("fontFamily", DEFAULT_THEME["fontFamily"])

    if chart_elem:
        chart_type = chart_elem.get("chart_type", "bar")
        categories = chart_elem.get("categories", [])
        series_data = chart_elem.get("series", [])

        # 映射图表类型
        chart_type_map = {
            "bar": XL_CHART_TYPE.BAR_CLUSTERED,
            "line": XL_CHART_TYPE.LINE,
            "pie": XL_CHART_TYPE.PIE,
            "scatter": XL_CHART_TYPE.SCATTER
        }
        xl_chart_type = chart_type_map.get(chart_type, XL_CHART_TYPE.BAR_CLUSTERED)

        # 创建图表数据
        from pptx.chart.data import CategoryChartData
        chart_data = CategoryChartData()
        chart_data.categories = categories

        for s in series_data:
            chart_data.add_series(s.get("name", ""), s.get("values", []))

        # 添加图表 — 16:9 适配：宽度从 Inches(8) 调整为 Inches(11)
        chart_frame = slide.shapes.add_chart(
            xl_chart_type,
            Inches(1), Inches(2), Inches(11), Inches(5),
            chart_data
        )

        chart = chart_frame.chart
        chart.has_legend = len(series_data) > 1
        if chart.has_legend:
            chart.legend.include_in_layout = False

    return slide


# ============================================================
# 幻灯片类型分发
# ============================================================

SLIDE_HANDLERS = {
    "title": add_title_slide,
    "content": add_content_slide,
    "two_content": add_two_content_slide,
    "table": add_table_slide,
    "image": add_image_slide,
    "section": add_section_slide,
    "blank": add_blank_slide,
    "chart": add_chart_slide,
}


# ============================================================
# 主逻辑
# ============================================================

def create_presentation(slides_data, theme, aspect_ratio="16:9", template_path=None):
    """创建新演示文稿

    若提供 template_path，基于模板创建（保留母版/主题色/字体）。
    否则使用默认空白模板。
    aspect_ratio 仅在无模板时生效（模板自带宽高比）。
    """
    if template_path:
        prs = Presentation(template_path)
    else:
        prs = Presentation()
        # 应用宽高比
        ratio = ASPECT_RATIOS.get(aspect_ratio, ASPECT_RATIOS["16:9"])
        prs.slide_width = ratio["width"]
        prs.slide_height = ratio["height"]

    for slide_data in slides_data:
        layout = slide_data.get("layout", "content")
        handler = SLIDE_HANDLERS.get(layout, add_content_slide)
        handler(prs, slide_data, theme)

    return prs


def modify_presentation(prs, slides_data, modify_action, target_slide, theme):
    """修改已有演示文稿

    保留源文件的宽高比，不覆盖。
    insert/replace 操作使用 try-except 保护内部 API 调用。
    """

    if modify_action == "append":
        for slide_data in slides_data:
            layout = slide_data.get("layout", "content")
            handler = SLIDE_HANDLERS.get(layout, add_content_slide)
            handler(prs, slide_data, theme)

    elif modify_action == "insert" and target_slide is not None:
        for i, slide_data in enumerate(slides_data):
            layout = slide_data.get("layout", "content")
            handler = SLIDE_HANDLERS.get(layout, add_content_slide)
            new_slide = handler(prs, slide_data, theme)
            slide_idx = len(prs.slides) - 1
            target_idx = target_slide - 1 + i
            if slide_idx != target_idx:
                try:
                    prs.slides._sldIdLst.remove(new_slide._element)
                    prs.slides._sldIdLst.insert(target_idx, new_slide._element)
                except Exception as e:
                    print(f"警告：插入幻灯片到位置 {target_idx + 1} 失败: {e}", file=sys.stderr)

    elif modify_action == "replace" and target_slide is not None:
        idx = target_slide - 1
        if 0 <= idx < len(prs.slides):
            # 删除旧幻灯片
            try:
                rId = prs.slides._sldIdLst[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]
            except Exception as e:
                print(f"警告：删除旧幻灯片失败: {e}", file=sys.stderr)

            # 添加新幻灯片并移动到目标位置
            for i, slide_data in enumerate(slides_data):
                layout = slide_data.get("layout", "content")
                handler = SLIDE_HANDLERS.get(layout, add_content_slide)
                new_slide = handler(prs, slide_data, theme)
                slide_idx = len(prs.slides) - 1
                target_idx = idx + i
                if slide_idx != target_idx:
                    try:
                        prs.slides._sldIdLst.remove(new_slide._element)
                        prs.slides._sldIdLst.insert(target_idx, new_slide._element)
                    except Exception as e:
                        print(f"警告：移动新幻灯片到位置 {target_idx + 1} 失败: {e}", file=sys.stderr)

    elif modify_action == "delete" and target_slide is not None:
        idx = target_slide - 1
        if 0 <= idx < len(prs.slides):
            try:
                rId = prs.slides._sldIdLst[idx].get('{http://schemas.openxmlformats.org/officeDocument/2006/relationships}id')
                prs.part.drop_rel(rId)
                del prs.slides._sldIdLst[idx]
            except Exception as e:
                print(f"警告：删除幻灯片失败: {e}", file=sys.stderr)

    return prs


def main():
    parser = argparse.ArgumentParser(description="PPT Creator - python-pptx 生成脚本")
    parser.add_argument("--input", required=True, help="slides JSON 文件路径")
    parser.add_argument("--output", required=True, help="输出 PPTX 文件路径")
    parser.add_argument("--skill-dir", help="技能目录绝对路径（优先于 __file__ 自探测）")
    parser.add_argument("--source", help="修改模式时，源 PPTX 文件路径")
    parser.add_argument("--aspect-ratio", default="16:9", choices=["16:9", "4:3"],
                        help="宽高比（仅 create 模式生效，默认 16:9）")
    parser.add_argument("--template", help="模板 PPTX 文件路径（仅 create 模式，基于模板创建）")
    args = parser.parse_args()

    # 读取 JSON
    with open(args.input, 'r', encoding='utf-8') as f:
        data = json.load(f)

    # JSON 校验
    validate_slides_json(data)

    title = data.get("title", "演示文稿")
    author = data.get("author", "")
    theme = data.get("theme", DEFAULT_THEME)
    mode = data.get("mode", "create")
    source = data.get("source") or args.source
    template = data.get("template") or args.template
    slides_data = data.get("slides", [])

    if mode == "create" or not source:
        # 创建新演示文稿
        prs = create_presentation(slides_data, theme, args.aspect_ratio, template)
    else:
        # 修改已有演示文稿
        if not os.path.isfile(source):
            print(f"错误：源文件不存在: {source}", file=sys.stderr)
            sys.exit(1)

        # 修改前自动备份（模块 8）
        backup_source(source, skill_dir)

        prs = Presentation(source)
        modify_action = data.get("modify_action", "append")
        target_slide = data.get("target_slide")
        prs = modify_presentation(prs, slides_data, modify_action, target_slide, theme)

    # 设置文档属性
    prs.core_properties.title = title
    if author:
        prs.core_properties.author = author

    # 保存
    prs.save(args.output)
    print(f"成功生成 PPTX: {args.output}")
    print(f"幻灯片数量: {len(prs.slides)}")


if __name__ == "__main__":
    main()